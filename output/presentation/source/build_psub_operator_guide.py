from __future__ import annotations

import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
PRESENTATION_DIR = ROOT / "output" / "presentation"
SOURCE_DIR = PRESENTATION_DIR / "source"
MANIFEST_PATH = SOURCE_DIR / "slide_manifest.json"
OUTPUT_PPTX = PRESENTATION_DIR / "PSUB-Operator-Guide.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLORS = {
    "navy": RGBColor(13, 27, 42),
    "slate": RGBColor(27, 38, 59),
    "steel": RGBColor(65, 90, 119),
    "blue": RGBColor(63, 114, 175),
    "gold": RGBColor(230, 190, 90),
    "mist": RGBColor(232, 238, 244),
    "text": RGBColor(231, 238, 247),
    "dark_text": RGBColor(18, 27, 38),
    "success": RGBColor(38, 166, 91),
    "warning": RGBColor(216, 119, 6),
    "danger": RGBColor(204, 51, 51),
}


def load_manifest() -> list[dict]:
    return json.loads(MANIFEST_PATH.read_text(encoding="utf-8"))


def add_full_bleed_background(slide, color=COLORS["navy"]):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_top_band(slide):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.32))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["blue"]
    band.line.fill.background()


def add_footer(slide, slide_number: int):
    footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(7.6), Inches(0.25))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "PSUB | Operator Guide"
    r.font.size = Pt(11)
    r.font.color.rgb = COLORS["mist"]
    p.alignment = PP_ALIGN.LEFT

    num = slide.shapes.add_textbox(Inches(12.2), Inches(7.0), Inches(0.5), Inches(0.25))
    tf = num.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = str(slide_number)
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = COLORS["mist"]
    p.alignment = PP_ALIGN.RIGHT


def add_title_block(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(8.9), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(24 if len(title) > 30 else 28)
    r.font.bold = True
    r.font.color.rgb = COLORS["text"]

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.28), Inches(8.4), Inches(0.5))
        tf = sub_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(13)
        r.font.color.rgb = COLORS["mist"]


def add_chip(slide, text: str, left, top, width, fill_color=COLORS["gold"], font_color=COLORS["dark_text"]):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.32))
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill_color
    chip.line.fill.background()
    tf = chip.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = font_color
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_bullets(slide, items: list[str], left, top, width, height, font_size=18, color=COLORS["text"]):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(5)
        p.line_spacing = 1.05
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return box


def add_callout_panel(slide, items: list[str], title="Operator Tips"):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.65), Inches(1.2), Inches(3.0), Inches(5.4))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["slate"]
    panel.line.color.rgb = COLORS["steel"]

    heading = slide.shapes.add_textbox(Inches(9.95), Inches(1.45), Inches(2.4), Inches(0.35))
    tf = heading.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(15)
    r.font.color.rgb = COLORS["gold"]

    add_bullets(slide, items, Inches(9.95), Inches(1.9), Inches(2.35), Inches(4.4), font_size=13)


def add_highlight_outline(slide, left, top, width, height, label: str):
    outline = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    outline.fill.background()
    outline.line.color.rgb = COLORS["gold"]
    outline.line.width = Pt(2.5)

    label_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top - 0.34),
        Inches(min(width, 2.95)),
        Inches(0.3),
    )
    label_box.fill.solid()
    label_box.fill.fore_color.rgb = COLORS["gold"]
    label_box.line.fill.background()
    tf = label_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = label
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = COLORS["dark_text"]
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_picture_cover(slide, image_path: Path, left, top, width, height):
    if not image_path.exists():
        raise FileNotFoundError(f"Missing image: {image_path}")
    with Image.open(image_path) as img:
        img_w, img_h = img.size
    target_ratio = width / height
    image_ratio = img_w / img_h

    pic = slide.shapes.add_picture(str(image_path), left, top, width=width if image_ratio >= target_ratio else None, height=height if image_ratio < target_ratio else None)
    if image_ratio > target_ratio:
        displayed_width = pic.width
        crop_total = (displayed_width - width) / displayed_width
        pic.crop_left = crop_total / 2
        pic.crop_right = crop_total / 2
        pic.left = left
        pic.top = top
        pic.width = width
        pic.height = height
    else:
        displayed_height = pic.height
        crop_total = (displayed_height - height) / displayed_height
        pic.crop_top = crop_total / 2
        pic.crop_bottom = crop_total / 2
        pic.left = left
        pic.top = top
        pic.width = width
        pic.height = height
    return pic


def add_screenshot_frame(slide, image_path: Path):
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.82), Inches(8.55), Inches(4.85))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(245, 248, 252)
    frame.line.color.rgb = COLORS["steel"]
    add_picture_cover(slide, image_path, Inches(0.73), Inches(2.0), Inches(8.2), Inches(4.5))


def add_proof_points(slide, items: list[dict[str, str]]):
    card_lefts = (Inches(0.72), Inches(3.38), Inches(6.04))
    for index, item in enumerate(items[:3]):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            card_lefts[index],
            Inches(2.05),
            Inches(2.35),
            Inches(1.4),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["slate"]
        card.line.color.rgb = COLORS["steel"]

        label_box = slide.shapes.add_textbox(card.left + Inches(0.22), card.top + Inches(0.16), Inches(1.75), Inches(0.28))
        tf = label_box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = item["label"]
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = COLORS["gold"]

        value_box = slide.shapes.add_textbox(card.left + Inches(0.22), card.top + Inches(0.46), Inches(1.75), Inches(0.66))
        tf = value_box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = item["value"]
        r.font.name = "Aptos Display"
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = COLORS["text"]
        p.alignment = PP_ALIGN.LEFT


def add_body_panel(slide, items: list[str], title="Key Points"):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.72),
        Inches(3.95),
        Inches(8.1),
        Inches(2.55),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["slate"]
    panel.line.color.rgb = COLORS["steel"]

    heading = slide.shapes.add_textbox(Inches(1.0), Inches(4.13), Inches(2.1), Inches(0.35))
    tf = heading.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(15)
    r.font.color.rgb = COLORS["gold"]

    body_font_size = 13 if len(items) > 3 else 14
    add_bullets(slide, items, Inches(1.0), Inches(4.5), Inches(7.72), Inches(1.7), font_size=body_font_size)


def render_hero_slide(prs, spec, slide_number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_path = (PRESENTATION_DIR / spec["image"]).resolve()
    add_picture_cover(slide, image_path, 0, 0, SLIDE_W, SLIDE_H)

    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = COLORS["navy"]
    overlay.fill.transparency = 0.28
    overlay.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(0.98), Inches(1.4), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["gold"]
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.72), Inches(0.72), Inches(5.6), Inches(0.95))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = spec["title"]
    r.font.name = "Aptos Display"
    r.font.size = Pt(27)
    r.font.bold = True
    r.font.color.rgb = COLORS["text"]

    subtitle_box = slide.shapes.add_textbox(Inches(0.78), Inches(1.68), Inches(5.5), Inches(0.65))
    tf = subtitle_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = spec.get("subtitle", "")
    r.font.name = "Aptos"
    r.font.size = Pt(13)
    r.font.color.rgb = COLORS["mist"]

    hero_shot = spec.get("hero_screenshot")
    if hero_shot:
        shot_shadow = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(6.45),
            Inches(1.18),
            Inches(6.1),
            Inches(4.55),
        )
        shot_shadow.fill.solid()
        shot_shadow.fill.fore_color.rgb = RGBColor(4, 10, 18)
        shot_shadow.fill.transparency = 0.42
        shot_shadow.line.fill.background()
        slide.shapes._spTree.remove(shot_shadow._element)
        slide.shapes._spTree.insert(2, shot_shadow._element)

        shot_frame = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(6.3),
            Inches(1.02),
            Inches(6.1),
            Inches(4.55),
        )
        shot_frame.fill.solid()
        shot_frame.fill.fore_color.rgb = RGBColor(248, 250, 252)
        shot_frame.line.color.rgb = COLORS["steel"]

        shot_header = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(6.55),
            Inches(1.24),
            Inches(1.1),
            Inches(0.28),
        )
        shot_header.fill.solid()
        shot_header.fill.fore_color.rgb = COLORS["gold"]
        shot_header.line.fill.background()
        tf = shot_header.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "Web UI"
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = COLORS["dark_text"]
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        shot_path = (PRESENTATION_DIR / hero_shot).resolve()
        add_picture_cover(slide, shot_path, Inches(6.48), Inches(1.46), Inches(5.72), Inches(3.82))

    add_chip(slide, spec["callouts"][0], Inches(0.74), Inches(5.9), Inches(3.55))
    add_footer(slide, slide_number)


def render_timeline_slide(prs, spec, slide_number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bleed_background(slide)
    add_top_band(slide)
    add_title_block(slide, spec["title"], spec.get("subtitle"))

    steps = spec.get("body", [])
    step_count = max(len(steps), 1)
    available_width = Inches(8.2)
    bubble_width = min(Inches(1.15), available_width / step_count - Inches(0.12))
    step_gap = (available_width - bubble_width) / max(step_count - 1, 1)
    start_left = Inches(0.78)
    line_top = Inches(3.15)
    connector_left = start_left + bubble_width / 2
    connector_width = step_gap * max(step_count - 1, 1)
    connector = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, connector_left, line_top, connector_width, Inches(0.06))
    connector.fill.solid()
    connector.fill.fore_color.rgb = COLORS["steel"]
    connector.line.fill.background()

    for idx, step in enumerate(steps, start=1):
        left = start_left + step_gap * (idx - 1)
        bubble = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(2.55), bubble_width, Inches(1.0))
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = COLORS["blue"] if idx % 2 else COLORS["steel"]
        bubble.line.fill.background()
        tf = bubble.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = str(idx)
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = COLORS["text"]
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        text = slide.shapes.add_textbox(left - Inches(0.12), Inches(3.72), bubble_width + Inches(0.24), Inches(1.05))
        tf = text.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(10.5)
            run.font.color.rgb = COLORS["text"]

    add_callout_panel(slide, spec.get("callouts", []))
    add_footer(slide, slide_number)


def render_checklist_slide(prs, spec, slide_number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bleed_background(slide)
    add_top_band(slide)
    add_title_block(slide, spec["title"], spec.get("subtitle"))

    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.95), Inches(8.55), Inches(4.9))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["slate"]
    panel.line.color.rgb = COLORS["steel"]

    for idx, item in enumerate(spec.get("body", [])):
        top = Inches(2.22) + Inches(0.67) * idx
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.92), top, Inches(0.28), Inches(0.28))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS["success"] if idx < 5 else COLORS["warning"]
        circle.line.fill.background()

        text = slide.shapes.add_textbox(Inches(1.28), top - Inches(0.04), Inches(7.45), Inches(0.38))
        tf = text.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = item
        for run in p.runs:
            run.font.size = Pt(15.5)
            run.font.color.rgb = COLORS["text"]

    add_callout_panel(slide, spec.get("callouts", []), title="Requirements")
    add_footer(slide, slide_number)


def render_image_focus_slide(prs, spec, slide_number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bleed_background(slide)
    add_top_band(slide)
    add_title_block(slide, spec["title"], spec.get("subtitle"))
    image_path = (PRESENTATION_DIR / spec["image"]).resolve()
    add_screenshot_frame(slide, image_path)
    for highlight in spec.get("highlights", []):
        add_highlight_outline(
            slide,
            highlight["left"],
            highlight["top"],
            highlight["width"],
            highlight["height"],
            highlight["label"],
        )
    add_callout_panel(slide, spec.get("body", []) + spec.get("callouts", []), title=spec.get("panel_title", "Operator Tips"))
    add_footer(slide, slide_number)


def render_two_column_slide(prs, spec, slide_number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bleed_background(slide)
    add_top_band(slide)
    add_title_block(slide, spec["title"], spec.get("subtitle"))

    left_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(1.95), Inches(5.85), Inches(4.95))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = COLORS["slate"]
    left_panel.line.color.rgb = COLORS["danger"]

    right_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.82), Inches(1.95), Inches(5.85), Inches(4.95))
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = COLORS["slate"]
    right_panel.line.color.rgb = COLORS["success"]

    add_chip(slide, "Common Failures", Inches(0.92), Inches(2.18), Inches(1.85), fill_color=COLORS["danger"], font_color=COLORS["text"])
    add_chip(slide, "Operator Actions", Inches(7.15), Inches(2.18), Inches(1.75), fill_color=COLORS["success"], font_color=COLORS["text"])
    add_bullets(slide, spec.get("body", []), Inches(0.95), Inches(2.65), Inches(5.1), Inches(3.9), font_size=15)
    add_bullets(slide, spec.get("callouts", []), Inches(7.15), Inches(2.65), Inches(5.1), Inches(3.9), font_size=15)
    add_footer(slide, slide_number)


def render_proof_points_slide(prs, spec, slide_number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bleed_background(slide)
    add_top_band(slide)
    add_title_block(slide, spec["title"], spec.get("subtitle"))
    add_proof_points(slide, spec.get("proof_points", []))
    add_body_panel(slide, spec.get("body", []), title=spec.get("body_title", "Key Points"))
    add_callout_panel(slide, spec.get("callouts", []), title=spec.get("panel_title", "Operator Angle"))
    add_footer(slide, slide_number)


def build_presentation():
    specs = load_manifest()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    renderers = {
        "hero": render_hero_slide,
        "timeline": render_timeline_slide,
        "checklist": render_checklist_slide,
        "image-focus": render_image_focus_slide,
        "two-column": render_two_column_slide,
        "proof-points": render_proof_points_slide,
    }

    for index, spec in enumerate(specs, start=1):
        layout = spec["layout"]
        if layout not in renderers:
            raise ValueError(f"Unsupported layout: {layout}")
        renderers[layout](prs, spec, index)

    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PPTX)
    print(f"Wrote {OUTPUT_PPTX}")


if __name__ == "__main__":
    try:
        build_presentation()
    except Exception as exc:
        print(f"Deck generation failed: {exc}", file=sys.stderr)
        raise
